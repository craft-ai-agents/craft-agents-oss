# /// script
# requires-python = ">=3.12"
# dependencies = [
#   "pymupdf>=1.27,<2",
#   "pypdf>=6.7,<7",
#   "click>=8.3,<9",
#   "openpyxl>=3.1,<4",
#   "python-pptx>=1.0,<2",
#   "python-docx>=1.1,<2",
# ]
# ///
"""PDF operations tool.

Commands:
  Organize:  extract, info, merge, split, rotate, reorder, duplicate
  Edit:      watermark, fill-form, compress, crop, resize, flatten, header-footer
  Security:  encrypt, decrypt, redact, sanitize
  Convert:   to-image, from-image, to-docx, to-xlsx, to-pptx

Usage:
    uv run pdf_tool.py COMMAND [OPTIONS]

Validation note:
    Page range parsing is strict. Out-of-range pages and malformed --pages input now
    return explicit errors instead of being silently ignored.
"""

import io
import json
import sys
from pathlib import Path

import click
import fitz  # pymupdf
from pypdf import PdfReader, PdfWriter
from pypdf.generic import BooleanObject, NameObject, RectangleObject

# Standard page sizes in points (width x height)
PAGE_SIZES: dict[str, tuple[float, float]] = {
    "a3": (841.89, 1190.55),
    "a4": (595.28, 841.89),
    "a5": (419.53, 595.28),
    "letter": (612, 792),
    "legal": (612, 1008),
    "tabloid": (792, 1224),
}


# --- Helpers ---


def write_output(text: str, output_path: str | None) -> None:
    """Write text to file or stdout."""
    if output_path:
        Path(output_path).write_text(text, encoding="utf-8")
        click.echo(f"Output written to {output_path}", err=True)
    else:
        click.echo(text)


def write_pdf(writer: PdfWriter, output_path: str) -> None:
    """Write a PdfWriter to a file."""
    with open(output_path, "wb") as f:
        writer.write(f)
    click.echo(f"PDF written to {output_path}", err=True)


def parse_page_range(range_str: str, total_pages: int) -> list[int]:
    """Parse a strict page range string like '1-3,5,7-9' into zero-based page indices.

    Validation is intentionally strict:
    - empty segments are invalid
    - pages must be in range 1..total_pages
    - range bounds must be ascending (start <= end)
    - invalid tokens fail fast (no silent dropping)
    """
    if total_pages <= 0:
        raise ValueError("PDF has no pages.")

    if not range_str or not range_str.strip():
        raise ValueError("Page range cannot be empty.")

    pages: list[int] = []
    seen: set[int] = set()
    parts = range_str.split(",")

    for raw_part in parts:
        part = raw_part.strip()
        if not part:
            raise ValueError("Invalid page range: empty segment detected.")

        if "-" in part:
            if part.count("-") != 1:
                raise ValueError(f"Invalid page range: '{part}'. Expected format: '1-3' or '5'.")
            start_s, end_s = part.split("-", 1)
            if not start_s.strip() or not end_s.strip():
                raise ValueError(f"Invalid page range: '{part}'. Expected format: '1-3' or '5'.")
            try:
                start = int(start_s)
                end = int(end_s)
            except ValueError:
                raise ValueError(f"Invalid page range: '{part}'. Expected format: '1-3' or '5'.")

            if start < 1 or end < 1:
                raise ValueError(f"Invalid page range: '{part}'. Pages must be >= 1.")
            if start > end:
                raise ValueError(f"Invalid page range: '{part}'. Start page must be <= end page.")
            if start > total_pages or end > total_pages:
                raise ValueError(f"Page range '{part}' out of bounds. Valid pages: 1-{total_pages}.")

            for page_num in range(start, end + 1):
                idx = page_num - 1
                if idx not in seen:
                    seen.add(idx)
                    pages.append(idx)
        else:
            try:
                page_num = int(part)
            except ValueError:
                raise ValueError(f"Invalid page number: '{part}'. Expected a number.")

            if page_num < 1 or page_num > total_pages:
                raise ValueError(f"Page '{page_num}' out of bounds. Valid pages: 1-{total_pages}.")

            idx = page_num - 1
            if idx not in seen:
                seen.add(idx)
                pages.append(idx)

    if not pages:
        raise ValueError("No pages selected.")

    return pages


def check_output_differs(file: str, output: str) -> None:
    """Ensure output file differs from input file."""
    if str(Path(file).resolve()) == str(Path(output).resolve()):
        click.echo("Error: output file cannot be the same as input file.", err=True)
        sys.exit(1)


def parse_color(color: str) -> tuple[float, float, float]:
    """Parse a color name or hex string to RGB tuple (0-1 range)."""
    color_map: dict[str, tuple[float, float, float]] = {
        "gray": (0.5, 0.5, 0.5), "grey": (0.5, 0.5, 0.5),
        "red": (1.0, 0.0, 0.0), "blue": (0.0, 0.0, 1.0),
        "green": (0.0, 0.5, 0.0), "black": (0.0, 0.0, 0.0),
        "white": (1.0, 1.0, 1.0),
    }
    if color.startswith("#") and len(color) == 7:
        try:
            r = int(color[1:3], 16) / 255.0
            g = int(color[3:5], 16) / 255.0
            b = int(color[5:7], 16) / 255.0
            return (r, g, b)
        except ValueError:
            pass
    return color_map.get(color.lower(), (0.5, 0.5, 0.5))


@click.group()
def cli() -> None:
    """PDF operations tool."""
    pass


# ============================================================
# Organize: extract, info, merge, split, rotate, reorder, duplicate
# ============================================================


@cli.command()
@click.argument("file", type=click.Path(exists=True, dir_okay=False))
@click.option("--pages", type=str, default=None, help="Page range to extract (e.g. '1-3,5,7-9'). 1-based.")
@click.option("-o", "--output", type=click.Path(), default=None, help="Write output to file.")
def extract(file: str, pages: str | None, output: str | None) -> None:
    """Extract text from a PDF file.

    Extracts all text by default, or specific pages with --pages.
    """
    try:
        doc = fitz.open(file)
        try:
            total = len(doc)

            if pages:
                page_indices = parse_page_range(pages, total)
            else:
                page_indices = list(range(total))

            parts: list[str] = []
            for idx in page_indices:
                page = doc[idx]
                text = page.get_text()
                parts.append(f"--- Page {idx + 1} ---\n{text}")
        finally:
            doc.close()
        write_output("\n".join(parts), output)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)


@cli.command()
@click.argument("file", type=click.Path(exists=True, dir_okay=False))
@click.option("-o", "--output", type=click.Path(), default=None, help="Write output to file.")
@click.option("--set", "set_mode", is_flag=True, default=False, help="Edit metadata (use with --title, --author, etc.).")
@click.option("--title", type=str, default=None, help="Set document title (requires --set).")
@click.option("--author", type=str, default=None, help="Set document author (requires --set).")
@click.option("--subject", type=str, default=None, help="Set document subject (requires --set).")
@click.option("--keywords", type=str, default=None, help="Set document keywords (requires --set).")
def info(file: str, output: str | None, set_mode: bool, title: str | None, author: str | None, subject: str | None, keywords: str | None) -> None:
    """Show or edit PDF metadata and information.

    Use --set with --title, --author, --subject, --keywords to edit metadata.
    """
    if set_mode:
        if not output:
            click.echo("Error: -o/--output is required when using --set.", err=True)
            sys.exit(1)
        check_output_differs(file, output)
        try:
            reader = PdfReader(file)
            writer = PdfWriter()
            writer.append(reader)
            meta: dict[str, str] = {}
            if title is not None:
                meta["/Title"] = title
            if author is not None:
                meta["/Author"] = author
            if subject is not None:
                meta["/Subject"] = subject
            if keywords is not None:
                meta["/Keywords"] = keywords
            if not meta:
                click.echo("Error: provide at least one of --title, --author, --subject, --keywords.", err=True)
                sys.exit(1)
            writer.add_metadata(meta)
            write_pdf(writer, output)
        except Exception as e:
            click.echo(f"Error: {e}", err=True)
            sys.exit(1)
        return

    # Read mode (existing behavior)
    try:
        reader = PdfReader(file)

        info_dict: dict[str, object] = {
            "file": str(Path(file).resolve()),
            "encrypted": reader.is_encrypted,
        }

        if reader.is_encrypted:
            info_dict["pages"] = None
            info_dict["note"] = "PDF is encrypted. Decrypt to view full metadata."
        else:
            info_dict["pages"] = len(reader.pages)
            meta_obj = reader.metadata

            if meta_obj:
                info_dict["metadata"] = {
                    "title": meta_obj.title,
                    "author": meta_obj.author,
                    "subject": meta_obj.subject,
                    "creator": meta_obj.creator,
                    "producer": meta_obj.producer,
                    "creation_date": str(meta_obj.creation_date) if meta_obj.creation_date else None,
                    "modification_date": str(meta_obj.modification_date) if meta_obj.modification_date else None,
                }

            # Page dimensions from first page
            if reader.pages:
                page = reader.pages[0]
                box = page.mediabox
                info_dict["page_size"] = {
                    "width": float(box.width),
                    "height": float(box.height),
                    "width_inches": round(float(box.width) / 72, 2),
                    "height_inches": round(float(box.height) / 72, 2),
                }

            # Form fields
            fields = reader.get_fields()
            if fields:
                field_info = []
                for name, field in fields.items():
                    field_info.append({
                        "name": name,
                        "type": str(field.get("/FT", "Unknown")),
                        "value": str(field.get("/V", "")),
                    })
                info_dict["form_fields"] = field_info

        write_output(json.dumps(info_dict, indent=2, default=str), output)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)


@cli.command()
@click.argument("files", nargs=-1, required=True, type=click.Path(exists=True, dir_okay=False))
@click.option("-o", "--output", type=click.Path(), required=True, help="Output PDF file path.")
def merge(files: tuple[str, ...], output: str) -> None:
    """Merge multiple PDF files into one.

    Files are merged in the order provided.
    """
    if len(files) < 2:
        click.echo("Error: at least 2 PDF files are required for merge.", err=True)
        sys.exit(1)

    resolved_output = str(Path(output).resolve())
    for f in files:
        if str(Path(f).resolve()) == resolved_output:
            click.echo("Error: output file cannot be one of the input files.", err=True)
            sys.exit(1)

    try:
        writer = PdfWriter()
        for f in files:
            reader = PdfReader(f)
            for page in reader.pages:
                writer.add_page(page)

        write_pdf(writer, output)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)


@cli.command()
@click.argument("file", type=click.Path(exists=True, dir_okay=False))
@click.option("--pages", type=str, required=True, help="Page range to extract (e.g. '1-3,5,7-9'). 1-based.")
@click.option("--exclude", is_flag=True, default=False, help="Exclude specified pages instead of extracting them.")
@click.option("-o", "--output", type=click.Path(), required=True, help="Output PDF file path.")
def split(file: str, pages: str, exclude: bool, output: str) -> None:
    """Split a PDF by extracting or excluding specific pages.

    Use --exclude to delete pages instead of extracting them.
    """
    check_output_differs(file, output)

    try:
        reader = PdfReader(file)
        total = len(reader.pages)
        page_indices = parse_page_range(pages, total)

        if exclude:
            page_indices = [i for i in range(total) if i not in set(page_indices)]

        if not page_indices:
            click.echo("Error: no valid pages in the specified range.", err=True)
            sys.exit(1)

        writer = PdfWriter()
        for idx in page_indices:
            writer.add_page(reader.pages[idx])

        write_pdf(writer, output)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)


@cli.command()
@click.argument("file", type=click.Path(exists=True, dir_okay=False))
@click.option("--angle", type=click.Choice(["90", "180", "270"]), required=True, help="Rotation angle in degrees.")
@click.option("--pages", type=str, default=None, help="Page range to rotate (default: all).")
@click.option("-o", "--output", type=click.Path(), required=True, help="Output PDF file path.")
def rotate(file: str, angle: str, pages: str | None, output: str) -> None:
    """Rotate PDF pages by 90, 180, or 270 degrees."""
    check_output_differs(file, output)
    try:
        reader = PdfReader(file)
        writer = PdfWriter()
        total = len(reader.pages)
        target = set(parse_page_range(pages, total)) if pages else set(range(total))
        for i, page in enumerate(reader.pages):
            if i in target:
                page.rotate(int(angle))
            writer.add_page(page)
        write_pdf(writer, output)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)


@cli.command()
@click.argument("file", type=click.Path(exists=True, dir_okay=False))
@click.option("--order", type=str, default=None, help="Comma-separated 1-based page order (e.g. '3,1,2,4').")
@click.option("--reverse", is_flag=True, default=False, help="Reverse page order.")
@click.option("-o", "--output", type=click.Path(), required=True, help="Output PDF file path.")
def reorder(file: str, order: str | None, reverse: bool, output: str) -> None:
    """Reorder PDF pages by custom order or reverse them."""
    check_output_differs(file, output)
    if order and reverse:
        click.echo("Error: --order and --reverse are mutually exclusive. Provide only one.", err=True)
        sys.exit(1)
    if not order and not reverse:
        click.echo("Error: provide --order or --reverse.", err=True)
        sys.exit(1)
    try:
        reader = PdfReader(file)
        total = len(reader.pages)
        if reverse:
            indices = list(reversed(range(total)))
        else:
            indices = [int(x.strip()) - 1 for x in order.split(",")]
            for idx in indices:
                if idx < 0 or idx >= total:
                    click.echo(f"Error: page {idx + 1} out of range (1-{total}).", err=True)
                    sys.exit(1)
        writer = PdfWriter()
        for idx in indices:
            writer.add_page(reader.pages[idx])
        write_pdf(writer, output)
    except ValueError:
        click.echo("Error: --order must be comma-separated page numbers.", err=True)
        sys.exit(1)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)


@cli.command()
@click.argument("file", type=click.Path(exists=True, dir_okay=False))
@click.option("--pages", type=str, required=True, help="Page range to duplicate (e.g. '1,3').")
@click.option("--copies", type=click.IntRange(min=2), default=2, help="Total copies of each page (min: 2, default: 2).")
@click.option("-o", "--output", type=click.Path(), required=True, help="Output PDF file path.")
def duplicate(file: str, pages: str, copies: int, output: str) -> None:
    """Duplicate specific pages within a PDF."""
    check_output_differs(file, output)
    try:
        reader = PdfReader(file)
        total = len(reader.pages)
        dup_set = set(parse_page_range(pages, total))
        writer = PdfWriter()
        for i, page in enumerate(reader.pages):
            writer.add_page(page)
            if i in dup_set:
                for _ in range(copies - 1):
                    writer.add_page(page)
        write_pdf(writer, output)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)


# ============================================================
# Edit: watermark, fill-form, compress, crop, resize, flatten, header-footer
# ============================================================


@cli.command()
@click.argument("file", type=click.Path(exists=True, dir_okay=False))
@click.option("--text", type=str, required=True, help="Watermark text.")
@click.option("--font-size", type=float, default=48, help="Font size (default: 48).")
@click.option("--opacity", type=float, default=0.3, help="Opacity 0-1 (default: 0.3).")
@click.option("--angle", type=float, default=45, help="Rotation angle in degrees (default: 45).")
@click.option("--color", type=str, default="gray", help="Color name or hex (default: gray).")
@click.option("-o", "--output", type=click.Path(), required=True, help="Output PDF file path.")
def watermark(file: str, text: str, font_size: float, opacity: float, angle: float, color: str, output: str) -> None:
    """Add a text watermark to every page of a PDF."""
    check_output_differs(file, output)

    try:
        doc = fitz.open(file)
        try:
            fitz_color = parse_color(color)

            for page in doc:
                rect = page.rect
                center = fitz.Point(rect.width / 2, rect.height / 2)

                # Measure text width to center it on the page
                text_width = fitz.get_text_length(text, fontsize=font_size)
                text_point = fitz.Point(center.x - text_width / 2, center.y + font_size / 3)

                tw = fitz.TextWriter(page.rect, opacity=opacity)
                tw.append(text_point, text, fontsize=font_size)
                rotation_matrix = fitz.Matrix(1, 1).prerotate(angle)
                tw.write_text(page, morph=(center, rotation_matrix), color=fitz_color)

            doc.save(output)
        finally:
            doc.close()
        click.echo(f"Watermarked PDF written to {output}", err=True)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)


@cli.command("fill-form")
@click.argument("file", type=click.Path(exists=True, dir_okay=False))
@click.option("--data", type=str, required=True, help="JSON string or path to JSON file with field values.")
@click.option("-o", "--output", type=click.Path(), required=True, help="Output PDF file path.")
def fill_form(file: str, data: str, output: str) -> None:
    """Fill PDF form fields with data from JSON.

    The JSON should map field names to values, e.g. {"name": "John", "date": "2024-01-01"}.
    """
    try:
        # Parse data - either JSON string or file path
        data_path = Path(data)
        if data_path.exists() and data_path.is_file():
            field_data = json.loads(data_path.read_text(encoding="utf-8"))
        else:
            field_data = json.loads(data)

        if not isinstance(field_data, dict):
            click.echo("Error: JSON data must be an object mapping field names to values.", err=True)
            sys.exit(1)

        reader = PdfReader(file)
        writer = PdfWriter()
        writer.append(reader)

        # Fill form fields
        for page_num in range(len(writer.pages)):
            writer.update_page_form_field_values(writer.pages[page_num], field_data)

        # Flatten if possible by setting NeedAppearances
        if "/AcroForm" in writer._root_object:
            writer._root_object["/AcroForm"][NameObject("/NeedAppearances")] = BooleanObject(True)

        write_pdf(writer, output)
    except json.JSONDecodeError as e:
        click.echo(f"Error parsing JSON data: {e}", err=True)
        sys.exit(1)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)


@cli.command()
@click.argument("file", type=click.Path(exists=True, dir_okay=False))
@click.option("-o", "--output", type=click.Path(), required=True, help="Output PDF file path.")
def compress(file: str, output: str) -> None:
    """Compress a PDF to reduce file size."""
    check_output_differs(file, output)
    try:
        doc = fitz.open(file)
        try:
            doc.save(output, deflate=True, garbage=4)
        finally:
            doc.close()
        original = Path(file).stat().st_size
        compressed = Path(output).stat().st_size
        ratio = (1 - compressed / original) * 100 if original > 0 else 0
        click.echo(
            f"Compressed: {original:,} -> {compressed:,} bytes ({ratio:.1f}% reduction). Written to {output}",
            err=True,
        )
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)


@cli.command()
@click.argument("file", type=click.Path(exists=True, dir_okay=False))
@click.option("--left", type=float, default=0, help="Points to trim from left.")
@click.option("--right", type=float, default=0, help="Points to trim from right.")
@click.option("--top", type=float, default=0, help="Points to trim from top.")
@click.option("--bottom", type=float, default=0, help="Points to trim from bottom.")
@click.option("--pages", type=str, default=None, help="Page range to crop (default: all).")
@click.option("-o", "--output", type=click.Path(), required=True, help="Output PDF file path.")
def crop(file: str, left: float, right: float, top: float, bottom: float, pages: str | None, output: str) -> None:
    """Crop PDF page margins by trimming points from each side."""
    check_output_differs(file, output)
    try:
        reader = PdfReader(file)
        writer = PdfWriter()
        total = len(reader.pages)
        target = set(parse_page_range(pages, total)) if pages else set(range(total))
        for i, page in enumerate(reader.pages):
            if i in target:
                box = page.mediabox
                page.cropbox = RectangleObject([
                    float(box.left) + left,
                    float(box.bottom) + bottom,
                    float(box.right) - right,
                    float(box.top) - top,
                ])
            writer.add_page(page)
        write_pdf(writer, output)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)


@cli.command()
@click.argument("file", type=click.Path(exists=True, dir_okay=False))
@click.option("--size", type=click.Choice(list(PAGE_SIZES.keys()), case_sensitive=False), required=True, help="Target page size.")
@click.option("--pages", type=str, default=None, help="Page range to resize (default: all).")
@click.option("-o", "--output", type=click.Path(), required=True, help="Output PDF file path.")
def resize(file: str, size: str, pages: str | None, output: str) -> None:
    """Resize PDF pages to a standard page size (A4, Letter, etc.)."""
    check_output_differs(file, output)
    target_w, target_h = PAGE_SIZES[size.lower()]
    try:
        reader = PdfReader(file)
        writer = PdfWriter()
        total = len(reader.pages)
        target_pages = set(parse_page_range(pages, total)) if pages else set(range(total))
        for i, page in enumerate(reader.pages):
            if i in target_pages:
                current_w = float(page.mediabox.width)
                current_h = float(page.mediabox.height)
                sx = target_w / current_w if current_w else 1
                sy = target_h / current_h if current_h else 1
                page.scale(float(sx), float(sy))
                page.mediabox = RectangleObject([0, 0, target_w, target_h])
            writer.add_page(page)
        write_pdf(writer, output)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)


@cli.command()
@click.argument("file", type=click.Path(exists=True, dir_okay=False))
@click.option("-o", "--output", type=click.Path(), required=True, help="Output PDF file path.")
def flatten(file: str, output: str) -> None:
    """Flatten PDF forms and annotations into static content."""
    check_output_differs(file, output)
    try:
        doc = fitz.open(file)
        try:
            new_doc = fitz.open()
            for page in doc:
                # show_pdf_page copies visual content (incl. annotations) as static content
                new_page = new_doc.new_page(width=page.rect.width, height=page.rect.height)
                new_page.show_pdf_page(new_page.rect, doc, page.number)
            new_doc.save(output)
            new_doc.close()
        finally:
            doc.close()
        click.echo(f"Flattened PDF written to {output}", err=True)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)


@cli.command("header-footer")
@click.argument("file", type=click.Path(exists=True, dir_okay=False))
@click.option("--header", "header_text", type=str, default=None, help="Header text. Use {n} for page number, {total} for total.")
@click.option("--footer", "footer_text", type=str, default=None, help="Footer text. Use {n} for page number, {total} for total.")
@click.option("--font-size", type=float, default=10, help="Font size (default: 10).")
@click.option("--color", type=str, default="black", help="Text color (default: black).")
@click.option("--margin", type=float, default=36, help="Margin from edge in points (default: 36 = 0.5in).")
@click.option("-o", "--output", type=click.Path(), required=True, help="Output PDF file path.")
def header_footer(file: str, header_text: str | None, footer_text: str | None, font_size: float, color: str, margin: float, output: str) -> None:
    """Add header and/or footer text to PDF pages.

    Use {n} for page number and {total} for total page count.
    Example: --footer "Page {n} of {total}"
    """
    check_output_differs(file, output)
    if not header_text and not footer_text:
        click.echo("Error: provide --header and/or --footer.", err=True)
        sys.exit(1)
    try:
        doc = fitz.open(file)
        try:
            total = len(doc)
            fitz_color = parse_color(color)
            for i, page in enumerate(doc):
                rect = page.rect
                page_num = i + 1
                if header_text:
                    text = header_text.replace("{n}", str(page_num)).replace("{total}", str(total))
                    text_width = fitz.get_text_length(text, fontsize=font_size)
                    x = (rect.width - text_width) / 2
                    y = margin + font_size
                    page.insert_text(fitz.Point(x, y), text, fontsize=font_size, color=fitz_color)
                if footer_text:
                    text = footer_text.replace("{n}", str(page_num)).replace("{total}", str(total))
                    text_width = fitz.get_text_length(text, fontsize=font_size)
                    x = (rect.width - text_width) / 2
                    y = rect.height - margin
                    page.insert_text(fitz.Point(x, y), text, fontsize=font_size, color=fitz_color)
            doc.save(output)
        finally:
            doc.close()
        click.echo(f"Header/footer added, written to {output}", err=True)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)


# ============================================================
# Security: encrypt, decrypt, redact, sanitize
# ============================================================


@cli.command()
@click.argument("file", type=click.Path(exists=True, dir_okay=False))
@click.option("--password", type=str, required=True, help="Password to set.")
@click.option("-o", "--output", type=click.Path(), required=True, help="Output PDF file path.")
def encrypt(file: str, password: str, output: str) -> None:
    """Password-protect a PDF with encryption."""
    check_output_differs(file, output)
    try:
        reader = PdfReader(file)
        writer = PdfWriter()
        writer.append(reader)
        writer.encrypt(user_password=password, owner_password=password)
        write_pdf(writer, output)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)


@cli.command()
@click.argument("file", type=click.Path(exists=True, dir_okay=False))
@click.option("--password", type=str, required=True, help="Password to unlock the PDF.")
@click.option("-o", "--output", type=click.Path(), required=True, help="Output PDF file path.")
def decrypt(file: str, password: str, output: str) -> None:
    """Remove password protection from an encrypted PDF."""
    check_output_differs(file, output)
    try:
        doc = fitz.open(file)
        try:
            if doc.is_encrypted:
                if not doc.authenticate(password):
                    click.echo("Error: incorrect password.", err=True)
                    sys.exit(1)
            doc.save(output)
        finally:
            doc.close()
        click.echo(f"Decrypted PDF written to {output}", err=True)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)


@cli.command()
@click.argument("file", type=click.Path(exists=True, dir_okay=False))
@click.option("--text", "texts", type=str, multiple=True, help="Text to redact (can specify multiple).")
@click.option("--area", "areas", type=str, multiple=True, help="Area to redact as 'page:x1,y1,x2,y2' (1-based page, points).")
@click.option("-o", "--output", type=click.Path(), required=True, help="Output PDF file path.")
def redact(file: str, texts: tuple[str, ...], areas: tuple[str, ...], output: str) -> None:
    """Permanently redact text or areas from a PDF.

    Redact by text:  --text "confidential" --text "SSN: 123-45-6789"
    Redact by area:  --area "1:100,200,300,250"  (page 1, rectangle in points)
    """
    check_output_differs(file, output)
    if not texts and not areas:
        click.echo("Error: provide --text and/or --area to redact.", err=True)
        sys.exit(1)
    try:
        doc = fitz.open(file)
        try:
            # Text-based redaction
            for search_text in texts:
                for page in doc:
                    instances = page.search_for(search_text)
                    for rect in instances:
                        page.add_redact_annot(rect, fill=(0, 0, 0))

            # Area-based redaction
            for area_spec in areas:
                try:
                    page_str, coords_str = area_spec.split(":", 1)
                    page_num = int(page_str) - 1
                    coords = [float(c.strip()) for c in coords_str.split(",")]
                    if len(coords) != 4:
                        raise ValueError("Need 4 coordinates")
                    rect = fitz.Rect(*coords)
                    doc[page_num].add_redact_annot(rect, fill=(0, 0, 0))
                except (ValueError, IndexError) as e:
                    click.echo(f"Error: invalid area spec '{area_spec}': {e}", err=True)
                    sys.exit(1)

            # Apply all redactions
            for page in doc:
                page.apply_redactions()

            doc.save(output)
        finally:
            doc.close()
        click.echo(f"Redacted PDF written to {output}", err=True)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)


@cli.command()
@click.argument("file", type=click.Path(exists=True, dir_okay=False))
@click.option("-o", "--output", type=click.Path(), required=True, help="Output PDF file path.")
def sanitize(file: str, output: str) -> None:
    """Clear metadata and perform aggressive PDF cleanup save.

    This command:
    - clears standard document metadata fields
    - removes XMP metadata
    - saves with garbage collection / cleanup options

    Note: it does not guarantee removal of every possible active-content pattern
    across all PDFs.
    """
    check_output_differs(file, output)
    try:
        doc = fitz.open(file)
        try:
            # Clear all standard metadata
            doc.set_metadata({
                "title": "",
                "author": "",
                "subject": "",
                "keywords": "",
                "creator": "",
                "producer": "",
            })
            # Remove XMP metadata
            doc.del_xml_metadata()
            # Save with full garbage collection to strip unreferenced objects (JS, embedded files, etc.)
            doc.save(output, garbage=4, clean=True, deflate=True)
        finally:
            doc.close()
        click.echo(f"Sanitized PDF written to {output}", err=True)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)


# ============================================================
# Convert: to-image, from-image, to-docx, to-xlsx, to-pptx
# ============================================================


@cli.command("to-image")
@click.argument("file", type=click.Path(exists=True, dir_okay=False))
@click.option("--format", "fmt", type=click.Choice(["png", "jpg"]), default="png", help="Image format (default: png).")
@click.option("--dpi", type=int, default=150, help="Resolution in DPI (default: 150).")
@click.option("--pages", type=str, default=None, help="Page range to convert (default: all).")
@click.option("-o", "--output", type=click.Path(), required=True, help="Output directory.")
def to_image(file: str, fmt: str, dpi: int, pages: str | None, output: str) -> None:
    """Convert PDF pages to images (PNG or JPG)."""
    try:
        doc = fitz.open(file)
        try:
            total = len(doc)
            indices = parse_page_range(pages, total) if pages else list(range(total))
            out_dir = Path(output)
            out_dir.mkdir(parents=True, exist_ok=True)
            for idx in indices:
                page = doc[idx]
                pix = page.get_pixmap(dpi=dpi)
                ext = "jpg" if fmt == "jpg" else "png"
                out_path = out_dir / f"page_{idx + 1}.{ext}"
                if fmt == "jpg":
                    pix.save(str(out_path), jpg_quality=85)
                else:
                    pix.save(str(out_path))
                click.echo(f"Saved {out_path}", err=True)
        finally:
            doc.close()
        click.echo(f"Converted {len(indices)} pages to {fmt.upper()} in {output}", err=True)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)


@cli.command("from-image")
@click.argument("files", nargs=-1, required=True, type=click.Path(exists=True, dir_okay=False))
@click.option("-o", "--output", type=click.Path(), required=True, help="Output PDF file path.")
def from_image(files: tuple[str, ...], output: str) -> None:
    """Convert images to a PDF document.

    Each image becomes one page. Supports PNG, JPG, BMP, GIF, SVG, WebP.
    """
    try:
        doc = fitz.open()
        for img_path in files:
            img = fitz.open(img_path)
            pdfbytes = img.convert_to_pdf()
            img_pdf = fitz.open("pdf", pdfbytes)
            doc.insert_pdf(img_pdf)
            img.close()
            img_pdf.close()
        doc.save(output)
        doc.close()
        click.echo(f"PDF written to {output} ({len(files)} pages from images)", err=True)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)


@cli.command("to-docx")
@click.argument("file", type=click.Path(exists=True, dir_okay=False))
@click.option("--pages", type=str, default=None, help="Page range to convert (default: all).")
@click.option("-o", "--output", type=click.Path(), required=True, help="Output DOCX file path.")
def to_docx(file: str, pages: str | None, output: str) -> None:
    """Convert PDF to Word document (text extraction based)."""
    from docx import Document as DocxDocument

    try:
        doc = fitz.open(file)
        try:
            total = len(doc)
            indices = parse_page_range(pages, total) if pages else list(range(total))
            docx_doc = DocxDocument()
            for i, idx in enumerate(indices):
                page = doc[idx]
                blocks = page.get_text("blocks")
                for block in sorted(blocks, key=lambda b: (b[1], b[0])):
                    if block[6] == 0:  # text block
                        text = block[4].strip()
                        if text:
                            docx_doc.add_paragraph(text)
                if i < len(indices) - 1:
                    docx_doc.add_page_break()
        finally:
            doc.close()
        docx_doc.save(output)
        click.echo(f"DOCX written to {output}", err=True)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)


@cli.command("to-xlsx")
@click.argument("file", type=click.Path(exists=True, dir_okay=False))
@click.option("--pages", type=str, default=None, help="Page range to convert (default: all).")
@click.option("-o", "--output", type=click.Path(), required=True, help="Output XLSX file path.")
def to_xlsx(file: str, pages: str | None, output: str) -> None:
    """Extract tables from PDF to Excel spreadsheet.

    Uses PyMuPDF's built-in table detection. Falls back to text extraction if no tables found.
    """
    from openpyxl import Workbook

    try:
        doc = fitz.open(file)
        try:
            total = len(doc)
            indices = parse_page_range(pages, total) if pages else list(range(total))
            wb = Workbook()
            wb.remove(wb.active)
            table_count = 0
            for idx in indices:
                page = doc[idx]
                tabs = page.find_tables()
                for j, table in enumerate(tabs.tables):
                    table_count += 1
                    ws = wb.create_sheet(title=f"Page{idx + 1}_Table{j + 1}")
                    for row in table.extract():
                        ws.append([cell if cell else "" for cell in row])
            if table_count == 0:
                # No tables found — extract text as fallback
                ws = wb.create_sheet(title="Text")
                for idx in indices:
                    page = doc[idx]
                    for line in page.get_text().split("\n"):
                        if line.strip():
                            ws.append([line.strip()])
        finally:
            doc.close()
        wb.save(output)
        click.echo(f"XLSX written to {output} ({table_count} tables extracted)", err=True)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)


@cli.command("to-pptx")
@click.argument("file", type=click.Path(exists=True, dir_okay=False))
@click.option("--dpi", type=int, default=200, help="Image resolution for slides (default: 200).")
@click.option("--pages", type=str, default=None, help="Page range to convert (default: all).")
@click.option("-o", "--output", type=click.Path(), required=True, help="Output PPTX file path.")
def to_pptx(file: str, dpi: int, pages: str | None, output: str) -> None:
    """Convert PDF pages to PowerPoint slides (image-based)."""
    from pptx import Presentation
    from pptx.util import Emu

    try:
        doc = fitz.open(file)
        try:
            total = len(doc)
            indices = parse_page_range(pages, total) if pages else list(range(total))
            if not indices:
                click.echo("Error: no pages selected for conversion.", err=True)
                sys.exit(1)
            prs = Presentation()
            # Set slide size to match first page aspect ratio (points -> EMU: 1pt = 12700 EMU)
            first_page = doc[indices[0]]
            prs.slide_width = Emu(int(first_page.rect.width * 12700))
            prs.slide_height = Emu(int(first_page.rect.height * 12700))
            blank_layout = prs.slide_layouts[6]  # blank slide
            for idx in indices:
                page = doc[idx]
                pix = page.get_pixmap(dpi=dpi)
                img_bytes = pix.tobytes("png")
                slide = prs.slides.add_slide(blank_layout)
                slide.shapes.add_picture(
                    io.BytesIO(img_bytes), 0, 0,
                    prs.slide_width, prs.slide_height,
                )
        finally:
            doc.close()
        prs.save(output)
        click.echo(f"PPTX written to {output} ({len(indices)} slides)", err=True)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)


if __name__ == "__main__":
    cli()
